import os
import logging
import base64
import requests
import tempfile
import json
import shutil
import qrcode
from barcode.codex import Code128, Code39
from barcode.ean import EAN13, EAN8
from barcode.upc import UPCA
from barcode.isxn import ISBN13, ISBN10, ISSN
from barcode.writer import ImageWriter
from datetime import datetime
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from flask import Flask, render_template, request, jsonify, session, send_file, make_response
from flask_sqlalchemy import SQLAlchemy
from sqlalchemy.orm import DeclarativeBase
from werkzeug.middleware.proxy_fix import ProxyFix
from openai_client import OpenAIClient

# 文件生成相關導入
try:
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import xlsxwriter
    HAS_XLSXWRITER = True
except ImportError:
    HAS_XLSXWRITER = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from docx import Document
    from docx.shared import Inches as DocxInches
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import gspread
    from google.auth.exceptions import GoogleAuthError
    HAS_GSPREAD = True
except ImportError:
    HAS_GSPREAD = False

# Configure logging
logging.basicConfig(level=logging.INFO)

# Disable HTTP and OpenAI debug logs
logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("openai").setLevel(logging.WARNING)
logging.getLogger("httpcore").setLevel(logging.WARNING)

class Base(DeclarativeBase):
    pass

app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "dev-secret-key")
app.wsgi_app = ProxyFix(app.wsgi_app, x_proto=1, x_host=1)

# Database configuration
app.config["SQLALCHEMY_DATABASE_URI"] = os.environ.get("DATABASE_URL")
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False
app.config["SQLALCHEMY_ENGINE_OPTIONS"] = {
    'pool_pre_ping': True,
    "pool_recycle": 300,
}

# Initialize database
db = SQLAlchemy(app, model_class=Base)

# Initialize OpenAI client
openai_client = OpenAIClient()

# Create database tables
with app.app_context():
    import models
    models.init_db(db)
    db.create_all()

@app.route('/')
def index():
    """Render the enhanced main chat interface"""
    return render_template('enhanced_index.html')

@app.route('/simple')
def simple_index():
    """Render the simple chat interface"""
    return render_template('index.html')

@app.route('/chat', methods=['POST'])
def chat():
    """Handle chat messages and return AI responses"""
    try:
        data = request.get_json()
        user_message = data.get('message', '').strip()
        
        if not user_message:
            return jsonify({'error': 'Message cannot be empty'}), 400
        
        # Get chat history from session - limit to prevent cookie overflow
        if 'chat_history' not in session:
            session['chat_history'] = []
        
        # Limit chat history to last 10 messages to prevent cookie size issues
        chat_history = session['chat_history']
        if len(chat_history) > 10:
            chat_history = chat_history[-10:]
            session['chat_history'] = chat_history
        
        # Add user message to history
        chat_history.append({
            'role': 'user',
            'content': user_message
        })
        
        # Check if this needs separate code responses
        needs_separation = any(keyword in user_message.lower() for keyword in [
            'apps script', 'app script', 'html', 'css', 'javascript', 'js', 
            'python', 'sql', 'php', 'java', 'c++', 'c#', 'react', 'vue', 'nodejs',
            'flutter', 'dart', 'swift', 'kotlin', 'go', 'rust', 'typescript',
            '程式碼', '前端', '後端', 'frontend', 'backend'
        ])
        
        # Also check if message contains multiple programming languages
        code_keywords = ['html', 'css', 'javascript', 'js', 'python', 'sql', 'php', 'java', 'c++', 'c#', 'csharp', 'react', 'vue', 'nodejs', 'node', 'flutter', 'dart', 'swift', 'kotlin', 'go', 'rust', 'typescript', 'ts']
        mentioned_languages = [lang for lang in code_keywords if lang in user_message.lower()]
        
        if needs_separation or len(mentioned_languages) > 1:
            # Determine code types based on context
            responses = []
            
            # Check for APPS SCRIPT
            if 'apps script' in user_message.lower() or 'app script' in user_message.lower():
                # Always provide GS Code
                gs_prompt = chat_history + [{
                    'role': 'user', 
                    'content': f"{user_message}\n\n請只提供正確的GS程式碼，要求：1. 使用標準Google Apps Script函式語法 2. 確保函式可以在Apps Script環境正常執行 3. 提供正確的函式調用方式"
                }]
                gs_response = openai_client.get_response(gs_prompt)
                responses.append(gs_response)
                
                # Only provide HTML if it's a UI-related request
                html_keywords = ['介面', '網頁', 'ui', '表單', '按鈕', '輸入', '顯示', '視窗', 'html', '前端', '使用者介面']
                needs_html = any(keyword in user_message.lower() for keyword in html_keywords)
                
                if needs_html:
                    html_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請只提供HTML程式碼部分，並包含使用教學：1. 如何在Apps Script中建立HTML檔案 2. 如何部署和使用"
                    }]
                    html_response = openai_client.get_response(html_prompt)
                    responses.append(html_response)
                
            # Check for web development and frameworks
            if any(lang in user_message.lower() for lang in ['html', 'css', 'javascript', 'js', 'react', 'vue', 'typescript', 'ts']):
                # HTML Code
                if 'html' in user_message.lower():
                    html_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的HTML程式碼，必須包含：\n1. 完整的文件結構和語義標籤\n2. 所有必要的meta標籤和屬性\n3. 響應式設計支援\n4. 可訪問性考量\n5. 實際可運行的完整HTML\n6. 詳細的註解說明\n7. 使用範例和部署教學\n8. CSS和JavaScript整合指引\n\n確保HTML程式碼符合現代標準且完整。"
                    }]
                    html_response = openai_client.get_response(html_prompt)
                    responses.append(html_response)
                
                # CSS Code
                if 'css' in user_message.lower():
                    css_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的CSS程式碼，必須包含：\n1. 完整的樣式結構和選擇器\n2. 響應式設計和媒體查詢\n3. 現代CSS特性和變數\n4. 瀏覽器兼容性考量\n5. 實際可使用的完整樣式\n6. 詳細的註解說明\n7. 使用範例和整合教學\n8. 效能優化建議\n\n確保CSS程式碼現代且高效。"
                    }]
                    css_response = openai_client.get_response(css_prompt)
                    responses.append(css_response)
                
                # React Code
                if 'react' in user_message.lower():
                    react_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供簡潔的React程式碼，包含：\n1. 核心組件結構\n2. 必要的import\n3. 基本狀態管理\n4. 事件處理\n5. 可運行的完整組件\n6. package.json設定\n\n回應簡潔，專注於核心功能。"
                    }]
                    react_response = openai_client.get_response(react_prompt)
                    responses.append(react_response)
                
                # Vue Code
                if 'vue' in user_message.lower():
                    vue_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的Vue.js程式碼，必須包含：\n1. 完整的組件結構和模板\n2. 所有必要的import和依賴\n3. 響應式數據和計算屬性\n4. 事件處理和指令使用\n5. 實際可運行的完整組件\n6. 詳細的註解說明\n7. 使用範例和部署教學\n8. Vue CLI和構建配置\n\n確保Vue程式碼符合框架最佳實踐。"
                    }]
                    vue_response = openai_client.get_response(vue_prompt)
                    responses.append(vue_response)
                
                # TypeScript Code
                if any(ts in user_message.lower() for ts in ['typescript', 'ts']):
                    typescript_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的TypeScript程式碼，必須包含：\n1. 完整的類型定義和介面\n2. 所有必要的import和類型宣告\n3. 錯誤處理和類型安全\n4. 完整的函數和類別定義\n5. 實際可編譯的完整程式碼\n6. 詳細的註解說明\n7. 使用範例和編譯教學\n8. tsconfig.json配置說明\n\n確保TypeScript程式碼類型安全且完整。"
                    }]
                    typescript_response = openai_client.get_response(typescript_prompt)
                    responses.append(typescript_response)
                
                # JavaScript Code
                if any(js in user_message.lower() for js in ['javascript', 'js']) and 'typescript' not in user_message.lower():
                    js_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供簡潔的JavaScript程式碼，包含：\n1. 核心程式邏輯\n2. 必要的變數和函數\n3. 基本錯誤處理\n4. DOM操作和事件\n5. 可執行的完整程式\n6. 簡要註解\n\n回應簡潔，專注於核心實現。"
                    }]
                    js_response = openai_client.get_response(js_prompt)
                    responses.append(js_response)
            
            # Check for backend and other languages
            if any(lang in user_message.lower() for lang in ['python', 'sql', 'php', 'java', 'c++', 'c#', 'csharp', 'nodejs', 'node', 'go', 'rust', 'kotlin', 'swift', 'dart', 'flutter']):
                # Python Code
                if 'python' in user_message.lower():
                    python_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供簡潔的Python程式碼，包含：\n1. 核心程式邏輯\n2. 必要的import\n3. 基本錯誤處理\n4. 可執行的完整程式\n5. 簡要註解\n6. 套件安裝指令\n\n回應簡潔，專注於核心實現。"
                    }]
                    python_response = openai_client.get_response(python_prompt)
                    responses.append(python_response)
                
                # Java Code
                if 'java' in user_message.lower():
                    java_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的Java程式碼，必須包含：\n1. 完整的類別結構和套件宣告\n2. 所有必要的import語句\n3. 異常處理機制\n4. 完整的方法定義和參數\n5. 實際可編譯執行的程式碼\n6. 詳細的註解說明\n7. 使用範例和編譯教學\n8. Maven/Gradle依賴配置\n\n確保Java程式碼符合最佳實踐且可直接使用。"
                    }]
                    java_response = openai_client.get_response(java_prompt)
                    responses.append(java_response)
                
                # C# Code
                if any(cs in user_message.lower() for cs in ['c#', 'csharp']):
                    csharp_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的C#程式碼，必須包含：\n1. 完整的命名空間和類別結構\n2. 所有必要的using語句\n3. 異常處理機制\n4. 完整的方法定義和屬性\n5. 實際可編譯執行的程式碼\n6. 詳細的註解說明\n7. 使用範例和執行教學\n8. NuGet套件依賴說明\n\n確保C#程式碼符合.NET最佳實踐。"
                    }]
                    csharp_response = openai_client.get_response(csharp_prompt)
                    responses.append(csharp_response)
                
                # Node.js Code
                if any(node in user_message.lower() for node in ['nodejs', 'node']):
                    nodejs_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供簡潔的Node.js程式碼，包含：\n1. 核心程式碼結構\n2. 必要的require語句\n3. 基本錯誤處理\n4. 可執行的完整程式\n5. 簡要註解\n6. package.json設定\n\n回應請保持簡潔，專注於核心功能實現。"
                    }]
                    nodejs_response = openai_client.get_response(nodejs_prompt)
                    responses.append(nodejs_response)
                
                # Flutter Code
                if 'flutter' in user_message.lower():
                    flutter_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供簡潔的Flutter程式碼，包含：\n1. 核心Widget結構\n2. 必要的import\n3. 基本狀態管理\n4. UI設計和互動\n5. 可執行的完整app\n6. pubspec.yaml設定\n\n回應簡潔，專注於核心功能。"
                    }]
                    flutter_response = openai_client.get_response(flutter_prompt)
                    responses.append(flutter_response)
                
                # Dart Code
                if 'dart' in user_message.lower() and 'flutter' not in user_message.lower():
                    dart_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的Dart程式碼，必須包含：\n1. 完整的Dart程式結構和類別\n2. 所有必要的import語句\n3. 錯誤處理和異常管理\n4. 完整的函數和方法定義\n5. 實際可執行的完整程式碼\n6. 詳細的註解說明\n7. 使用範例和執行教學\n8. pubspec.yaml配置說明\n\n確保Dart程式碼現代且高效。"
                    }]
                    dart_response = openai_client.get_response(dart_prompt)
                    responses.append(dart_response)
                
                # Swift Code
                if 'swift' in user_message.lower():
                    swift_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的Swift程式碼，必須包含：\n1. 完整的Swift程式結構和類別\n2. 所有必要的import框架\n3. 錯誤處理和可選型使用\n4. 完整的函數和方法定義\n5. 實際可編譯的iOS應用程式\n6. 詳細的註解說明\n7. 使用範例和Xcode配置\n8. Info.plist和權限設定\n\n確保Swift程式碼符合iOS開發最佳實踐。"
                    }]
                    swift_response = openai_client.get_response(swift_prompt)
                    responses.append(swift_response)
                
                # Kotlin Code
                if 'kotlin' in user_message.lower():
                    kotlin_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的Kotlin程式碼，必須包含：\n1. 完整的Kotlin程式結構和類別\n2. 所有必要的import語句\n3. 錯誤處理和Null安全\n4. 完整的函數和擴展定義\n5. 實際可編譯的Android應用程式\n6. 詳細的註解說明\n7. 使用範例和Android Studio配置\n8. build.gradle依賴設定\n\n確保Kotlin程式碼符合Android開發最佳實踐。"
                    }]
                    kotlin_response = openai_client.get_response(kotlin_prompt)
                    responses.append(kotlin_response)
                
                # Go Code
                if 'go' in user_message.lower() and 'golang' not in user_message.lower():
                    go_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的Go程式碼，必須包含：\n1. 完整的Go程式結構和套件\n2. 所有必要的import語句\n3. 錯誤處理和goroutine\n4. 完整的函數和結構定義\n5. 實際可編譯執行的程式碼\n6. 詳細的註解說明\n7. 使用範例和模組管理\n8. go.mod依賴配置\n\n確保Go程式碼高效且符合慣例。"
                    }]
                    go_response = openai_client.get_response(go_prompt)
                    responses.append(go_response)
                
                # Rust Code
                if 'rust' in user_message.lower():
                    rust_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的Rust程式碼，必須包含：\n1. 完整的Rust程式結構和模組\n2. 所有必要的use語句和crate\n3. 記憶體安全和所有權管理\n4. 完整的函數和結構定義\n5. 實際可編譯執行的程式碼\n6. 詳細的註解說明\n7. 使用範例和Cargo配置\n8. Cargo.toml依賴設定\n\n確保Rust程式碼安全且高效。"
                    }]
                    rust_response = openai_client.get_response(rust_prompt)
                    responses.append(rust_response)
                
                # SQL Code
                if 'sql' in user_message.lower():
                    sql_prompt = chat_history + [{
                        'role': 'user', 
                        'content': f"{user_message}\n\n請提供完整的SQL程式碼，必須包含：\n1. 完整的SQL語句和語法\n2. 必要的資料表結構定義\n3. 索引建議和效能優化\n4. 錯誤處理和約束條件\n5. 實際可執行的查詢語句\n6. 詳細的註解說明\n7. 使用範例和執行教學\n8. 不同資料庫系統的兼容性說明\n\n確保SQL程式碼完整、高效且可直接執行。"
                    }]
                    sql_response = openai_client.get_response(sql_prompt)
                    responses.append(sql_response)
            
            # Add all responses to history
            for response in responses:
                chat_history.append({'role': 'assistant', 'content': response})
            
            # Update session
            session['chat_history'] = chat_history
            session.modified = True
            
            return jsonify({
                'responses': responses,
                'is_multi_code': True,
                'success': True
            })
        else:
            # 檢查是否為圖片生成請求
            image_keywords = ['畫', '繪', '圖', 'draw', 'paint', 'image', 'picture', '生成圖片', '產圖', '畫一個', '畫個', '給我一個', '圖片']
            is_image_request = any(keyword in user_message.lower() for keyword in image_keywords)
            
            if is_image_request and openai_client.is_configured():
                try:
                    # 生成圖片
                    response = openai_client.client.images.generate(
                        model="dall-e-3",
                        prompt=user_message,
                        size="1024x1024",
                        quality="standard",
                        n=1,
                    )
                    
                    image_url = response.data[0].url
                    
                    # 下載圖片並轉換為 base64
                    img_response = requests.get(image_url)
                    if img_response.status_code == 200:
                        img = Image.open(BytesIO(img_response.content))
                        
                        # 轉換為 base64
                        buffer = BytesIO()
                        img.save(buffer, format='PNG')
                        img_base64 = base64.b64encode(buffer.getvalue()).decode()
                        
                        ai_response = f"我為您生成了圖片：{user_message}"
                        
                        # Add both user and AI response to history
                        chat_history.append({
                            'role': 'assistant',
                            'content': ai_response
                        })
                        
                        # Update session
                        session['chat_history'] = chat_history
                        session.modified = True
                        
                        return jsonify({
                            'response': ai_response,
                            'image_generated': True,
                            'image_url': f"data:image/png;base64,{img_base64}",
                            'success': True
                        })
                    else:
                        # 圖片下載失敗，返回普通回應
                        ai_response = openai_client.get_response(chat_history)
                except Exception as e:
                    logging.error(f"圖片生成失敗: {e}")
                    # 生成失敗，返回普通回應
                    ai_response = openai_client.get_response(chat_history)
            else:
                # 檢查是否需要網路搜尋（天氣、新聞等）
                if openai_client.needs_web_search(user_message):
                    search_type = openai_client.get_search_type(user_message)
                    search_results = openai_client.perform_web_search(user_message, search_type)
                    
                    if search_results:
                        # 將搜尋結果格式化並整合到回應中
                        formatted_results = openai_client.format_search_results(search_results, search_type)
                        
                        # 使用AI結合搜尋結果生成回應
                        enhanced_prompt = chat_history + [{
                            'role': 'user',
                            'content': f"{user_message}\n\n最新資訊：{formatted_results}\n\n請根據上述最新資訊回答用戶的問題。"
                        }]
                        ai_response = openai_client.get_response(enhanced_prompt)
                    else:
                        # 搜尋失敗，使用一般回應
                        ai_response = openai_client.get_response(chat_history)
                else:
                    # Regular single response
                    ai_response = openai_client.get_response(chat_history)
            
            # Add AI response to history
            chat_history.append({
                'role': 'assistant',
                'content': ai_response
            })
            
            # Update session with limited history
            session['chat_history'] = chat_history
            session.modified = True
            
            return jsonify({
                'response': ai_response,
                'success': True
            })
        
    except Exception as e:
        logging.error(f"Error in chat endpoint: {str(e)}")
        return jsonify({
            'error': f'AI service error: {str(e)}',
            'success': False
        }), 500

@app.route('/generate_image', methods=['POST'])
def generate_image():
    """Generate image using AI"""
    try:
        data = request.get_json()
        prompt = data.get('prompt', '').strip()
        
        if not prompt:
            return jsonify({'success': False, 'error': '請提供圖片描述'})
        
        if not openai_client.is_configured():
            return jsonify({'success': False, 'error': 'OpenAI API 未配置'})
        
        # 使用 OpenAI DALL-E 生成圖片
        response = openai_client.client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        
        image_url = response.data[0].url
        
        # 下載圖片並轉換為 base64
        img_response = requests.get(image_url)
        if img_response.status_code == 200:
            img = Image.open(BytesIO(img_response.content))
            
            # 轉換為 base64
            buffer = BytesIO()
            img.save(buffer, format='PNG')
            img_base64 = base64.b64encode(buffer.getvalue()).decode()
            
            return jsonify({
                'success': True,
                'image_url': f"data:image/png;base64,{img_base64}",
                'original_url': image_url,
                'prompt': prompt
            })
        else:
            return jsonify({'success': False, 'error': '無法下載生成的圖片'})
            
    except Exception as e:
        logging.error(f"圖片生成錯誤: {e}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/analyze_image', methods=['POST'])
def analyze_image():
    """Analyze uploaded image using AI"""
    try:
        data = request.get_json()
        image_base64 = data.get('image', '').strip()
        
        if not image_base64:
            return jsonify({'success': False, 'error': '請提供圖片'})
        
        if not openai_client.is_configured():
            return jsonify({'success': False, 'error': 'OpenAI API 未配置'})
        
        # 使用 OpenAI GPT-4o 分析圖片
        response = openai_client.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "請詳細分析這張圖片，描述其內容、風格、顏色、構圖等元素。"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{image_base64}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=500
        )
        
        analysis = response.choices[0].message.content
        
        return jsonify({
            'success': True,
            'analysis': analysis
        })
        
    except Exception as e:
        logging.error(f"圖片分析錯誤: {e}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/voice_to_text', methods=['POST'])
def voice_to_text():
    """Convert voice to text using AI"""
    try:
        if 'audio' not in request.files:
            return jsonify({'success': False, 'error': '請提供音頻檔案'})
        
        audio_file = request.files['audio']
        
        if not openai_client.is_configured():
            return jsonify({'success': False, 'error': 'OpenAI API 未配置'})
        
        # 使用 OpenAI Whisper 轉錄音頻
        transcript = openai_client.client.audio.transcriptions.create(
            model="whisper-1",
            file=audio_file
        )
        
        return jsonify({
            'success': True,
            'text': transcript.text
        })
        
    except Exception as e:
        logging.error(f"語音轉文字錯誤: {e}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/text_to_speech', methods=['POST'])
def text_to_speech():
    """Convert text to speech using AI"""
    try:
        data = request.get_json()
        text = data.get('text', '').strip()
        voice = data.get('voice', 'alloy')  # alloy, echo, fable, onyx, nova, shimmer
        
        if not text:
            return jsonify({'success': False, 'error': '請提供文字'})
        
        if not openai_client.is_configured():
            return jsonify({'success': False, 'error': 'OpenAI API 未配置'})
        
        # 使用 OpenAI TTS 生成語音
        response = openai_client.client.audio.speech.create(
            model="tts-1",
            voice=voice,
            input=text,
        )
        
        # 將音頻轉換為 base64
        audio_base64 = base64.b64encode(response.content).decode()
        
        return jsonify({
            'success': True,
            'audio_url': f"data:audio/mp3;base64,{audio_base64}"
        })
        
    except Exception as e:
        logging.error(f"文字轉語音錯誤: {e}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/save_game_score', methods=['POST'])
def save_game_score():
    """Save game score"""
    try:
        data = request.get_json()
        game_name = data.get('game_name')
        score = data.get('score')
        details = data.get('details', {})
        
        # 初始化遊戲分數記錄
        if 'game_scores' not in session:
            session['game_scores'] = {}
        
        if game_name not in session['game_scores']:
            session['game_scores'][game_name] = []
        
        # 記錄分數
        score_record = {
            'score': score,
            'details': details,
            'timestamp': datetime.now().isoformat(),
            'date': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
        
        session['game_scores'][game_name].append(score_record)
        
        # 只保留最近10次記錄
        session['game_scores'][game_name] = session['game_scores'][game_name][-10:]
        
        return jsonify({
            'success': True,
            'message': '分數已記錄',
            'best_score': max([r['score'] for r in session['game_scores'][game_name]])
        })
        
    except Exception as e:
        logging.error(f"Error saving game score: {str(e)}")
        return jsonify({'error': 'Failed to save score', 'success': False}), 500

@app.route('/get_game_scores', methods=['GET'])
def get_game_scores():
    """Get game scores"""
    try:
        game_name = request.args.get('game_name')
        
        if 'game_scores' not in session:
            session['game_scores'] = {}
        
        if game_name:
            scores = session['game_scores'].get(game_name, [])
            best_score = max([r['score'] for r in scores]) if scores else 0
            return jsonify({
                'success': True,
                'scores': scores,
                'best_score': best_score,
                'total_games': len(scores)
            })
        else:
            # 返回所有遊戲統計
            all_stats = {}
            for game, scores in session['game_scores'].items():
                all_stats[game] = {
                    'best_score': max([r['score'] for r in scores]) if scores else 0,
                    'total_games': len(scores),
                    'latest_score': scores[-1]['score'] if scores else 0
                }
            
            return jsonify({
                'success': True,
                'game_stats': all_stats
            })
        
    except Exception as e:
        logging.error(f"Error getting game scores: {str(e)}")
        return jsonify({'error': 'Failed to get scores', 'success': False}), 500

@app.route('/generate_document', methods=['POST'])
def generate_document():
    """Generate documents (Excel, Google Sheets, Word)"""
    try:
        data = request.get_json()
        method = data.get('method', 'ai_generate')
        
        if method == 'ai_generate':
            return generate_ai_document(data)
        elif method == 'template':
            return generate_template_document(data)
        elif method == 'ppt':
            return generate_ppt_document(data)
        else:
            return jsonify({'success': False, 'error': '無效的生成方法'})
            
    except Exception as e:
        logging.error(f"文件生成錯誤: {e}")
        return jsonify({'success': False, 'error': str(e)})

def generate_ai_document(data):
    """使用 AI 生成文件"""
    description = data.get('description', '')
    doc_type = data.get('type', 'excel')
    language = data.get('language', 'zh-TW')
    
    if not openai_client.is_configured():
        return jsonify({'success': False, 'error': 'OpenAI API 未配置'})
    
    # 使用 AI 生成文件結構
    prompt = f"""
    請根據以下描述生成一個{doc_type}文件的結構：
    描述：{description}
    語言：{language}
    
    請返回JSON格式，包含以下內容：
    - title: 文件標題
    - description: 文件說明
    - structure: 文件結構（對於試算表是columns和rows，對於文件是sections）
    
    確保內容實用且符合需求。
    """
    
    try:
        response = openai_client.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "你是一個專業的文件生成助手，擅長創建各種商業和學術文件。"},
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"}
        )
        
        ai_structure = json.loads(response.choices[0].message.content)
        
        # 根據類型生成對應文件
        if doc_type in ['excel', 'google_sheet']:
            return create_ai_excel_document(ai_structure, [])
        elif doc_type == 'word':
            return create_ai_word_document(ai_structure, [])
            
    except Exception as e:
        logging.error(f"AI 生成錯誤: {e}")
        return jsonify({'success': False, 'error': f'AI 生成失敗: {str(e)}'})

def generate_template_document(data):
    """使用範本生成文件"""
    template_type = data.get('template_type')
    template_name = data.get('template_name')
    language = data.get('language', 'zh-TW')
    
    templates = get_document_templates(language)
    
    if template_type == 'spreadsheet':
        if 'spreadsheet' in templates and template_name in templates['spreadsheet']:
            template_data = templates['spreadsheet'][template_name]
            return create_spreadsheet(template_data, 'excel', language)
    elif template_type == 'document':
        if 'document' in templates and template_name in templates['document']:
            template_data = templates['document'][template_name]
            return create_word_document(template_data, language)
    elif template_type == 'google_sheets':
        if 'google_sheets' in templates and template_name in templates['google_sheets']:
            template_data = templates['google_sheets'][template_name]
            return create_google_sheets_template(template_data, language)
    
    # 提供更詳細的錯誤訊息
    available_templates = []
    if 'spreadsheet' in templates:
        available_templates.extend([f"spreadsheet/{k}" for k in templates['spreadsheet'].keys()])
    if 'document' in templates:
        available_templates.extend([f"document/{k}" for k in templates['document'].keys()])
    if 'google_sheets' in templates:
        available_templates.extend([f"google_sheets/{k}" for k in templates['google_sheets'].keys()])
    
    error_msg = f'找不到範本 "{template_type}/{template_name}"。可用範本：{", ".join(available_templates)}'
    return jsonify({'success': False, 'error': error_msg})

@app.route('/generate_ai_document', methods=['POST'])
def generate_ai_document_endpoint():
    """AI智能文件生成端點"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({'success': False, 'error': '沒有提供數據'})
        
        document_type = data.get('document_type', 'excel')
        description = data.get('description', '')
        style = data.get('style', 'professional')
        language = data.get('language', 'zh-TW')
        include_images = data.get('include_images', False)
        
        if not description:
            return jsonify({'success': False, 'error': '請提供文件描述'})
        
        # 生成基本文件結構
        ai_structure = {
            "title": f"{description[:50]}..." if len(description) > 50 else description,
            "description": f"根據「{description}」生成的{document_type}文件",
            "structure": get_default_structure(document_type),
            "image_requirements": []
        }
        
        # 為Word和PPT生成圖片
        generated_images = []
        if include_images and document_type in ['word', 'ppt'] and openai_client.is_configured():
            try:
                # 根據描述生成相關圖片
                img_prompts = [
                    f"{description}, professional style",
                    f"diagram related to {description}, business style"
                ]
                
                for i, prompt in enumerate(img_prompts[:2]):  # 最多2張圖片
                    try:
                        img_response = openai_client.client.images.generate(
                            model="dall-e-3",
                            prompt=prompt[:1000],  # 限制提示長度
                            size="1024x1024",
                            quality="standard",
                            n=1
                        )
                        
                        # 下載並保存圖片
                        img_url = img_response.data[0].url
                        import requests
                        img_data = requests.get(img_url, timeout=30).content
                        
                        # 創建唯一檔名
                        import uuid
                        img_filename = f"ai_image_{uuid.uuid4()}.png"
                        img_path = os.path.join('static', 'downloads', img_filename)
                        os.makedirs(os.path.dirname(img_path), exist_ok=True)
                        
                        with open(img_path, 'wb') as img_file:
                            img_file.write(img_data)
                        
                        generated_images.append({
                            'url': f'/static/downloads/{img_filename}',
                            'description': f"圖片 {i+1}",
                            'path': img_path
                        })
                        
                    except Exception as e:
                        logging.warning(f"圖片生成失敗: {e}")
                        continue
                        
            except Exception as e:
                logging.warning(f"整體圖片生成過程失敗: {e}")
        
        # 根據文件類型生成相應文件
        if document_type == 'excel':
            return create_ai_excel_document(ai_structure, [])
        elif document_type == 'word':
            return create_ai_word_document(ai_structure, generated_images)
        elif document_type == 'ppt':
            return create_ai_ppt_document(ai_structure, generated_images)
        else:
            return jsonify({'success': False, 'error': '不支援的文件類型'})
        
    except Exception as e:
        logging.error(f"AI文件生成錯誤: {e}")
        return jsonify({'success': False, 'error': f'AI文件生成失敗: {str(e)}'})

def get_default_structure(document_type):
    """獲取預設文件結構"""
    if document_type == 'excel':
        return {
            "headers": ["項目", "描述", "數值", "備註"],
            "data": [
                ["項目1", "說明1", "100", ""],
                ["項目2", "說明2", "200", ""],
                ["項目3", "說明3", "300", ""]
            ]
        }
    elif document_type == 'word':
        return {
            "content": "這是AI生成的Word文件內容。本文件根據用戶需求自動創建，包含相關的段落和格式。"
        }
    elif document_type == 'ppt':
        return {
            "slides": [
                {"title": "標題頁", "content": "這是AI生成的簡報"},
                {"title": "內容概述", "content": "• 重點一\n• 重點二\n• 重點三"},
                {"title": "結論", "content": "總結與展望"}
            ]
        }
    else:
        return {}

def generate_ai_document_with_images(document_type, description, style, language, include_images):
    """使用AI生成帶有圖片的文件"""
    try:
        # 生成文件結構
        structure_prompt = f"""
        請根據以下要求生成一個{document_type}文件的詳細結構：
        
        文件描述：{description}
        風格：{style}
        語言：{language}
        
        請以JSON格式回應，包含：
        {{
            "title": "文件標題",
            "description": "文件描述",
            "structure": {{
                "headers": ["欄位1", "欄位2", "欄位3"] (適用於Excel),
                "content": "文件內容段落" (適用於Word),
                "slides": [
                    {{
                        "title": "投影片標題",
                        "content": "投影片內容"
                    }}
                ] (適用於PPT)
            }},
            "image_requirements": ["圖片描述1", "圖片描述2"] (如需圖片，最多2個)
        }}
        """
        
        if not openai_client.client:
            return jsonify({'success': False, 'error': 'OpenAI客戶端未初始化'})
        
        response = openai_client.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "你是一個專業的文件生成專家，能夠創建結構化的商業文件。請確保回應是有效的JSON格式。"},
                {"role": "user", "content": structure_prompt}
            ],
            response_format={"type": "json_object"},
            max_tokens=2000
        )
        
        try:
            ai_structure = json.loads(response.choices[0].message.content)
        except json.JSONDecodeError as e:
            logging.error(f"JSON解析錯誤: {e}")
            # 提供備用結構
            ai_structure = {
                "title": f"AI生成的{document_type}文件",
                "description": description[:100] + "...",
                "structure": get_default_structure(document_type),
                "image_requirements": ["professional business image", "relevant diagram"] if include_images else []
            }
        
        # 生成圖片（如果需要）
        generated_images = []
        if include_images and 'image_requirements' in ai_structure and document_type in ['word', 'ppt']:
            for img_desc in ai_structure['image_requirements'][:2]:  # 最多生成2張圖片避免超時
                try:
                    # 簡化圖片描述避免過長
                    simplified_desc = img_desc[:100] if len(img_desc) > 100 else img_desc
                    
                    img_response = openai_client.client.images.generate(
                        model="dall-e-3",
                        prompt=f"{simplified_desc}, {style} style, high quality",
                        size="1024x1024",
                        quality="standard",
                        n=1
                    )
                    
                    # 下載並保存圖片
                    img_url = img_response.data[0].url
                    
                    # 使用timeout避免長時間等待
                    import requests
                    img_data = requests.get(img_url, timeout=30).content
                    
                    # 創建唯一檔名
                    import uuid
                    img_filename = f"ai_image_{uuid.uuid4()}.png"
                    img_path = os.path.join('static', 'downloads', img_filename)
                    os.makedirs(os.path.dirname(img_path), exist_ok=True)
                    
                    with open(img_path, 'wb') as img_file:
                        img_file.write(img_data)
                    
                    generated_images.append({
                        'url': f'/static/downloads/{img_filename}',
                        'description': img_desc,
                        'path': img_path
                    })
                    
                except Exception as e:
                    logging.warning(f"圖片生成失敗: {e}")
                    # 繼續處理其他圖片，不中斷整個流程
                    continue
        
        # 根據文件類型生成相應文件
        if document_type == 'excel':
            return create_ai_excel_document(ai_structure, generated_images)
        elif document_type == 'word':
            return create_ai_word_document(ai_structure, generated_images)
        elif document_type == 'ppt':
            return create_ai_ppt_document(ai_structure, generated_images)
        else:
            return jsonify({'success': False, 'error': '不支援的文件類型'})
            
    except Exception as e:
        logging.error(f"AI文件生成錯誤: {e}")
        return jsonify({'success': False, 'error': f'文件生成失敗: {str(e)}'})

def create_ai_excel_document(structure, images):
    """創建AI生成的Excel文件"""
    try:
        if not HAS_OPENPYXL:
            return jsonify({'success': False, 'error': '系統未安裝 openpyxl 套件'})
        
        # 創建工作簿
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = structure.get('title', 'AI生成表格')
        
        # 設置樣式
        title_font = Font(size=16, bold=True)
        header_font = Font(size=12, bold=True)
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        
        # 添加標題
        title = structure.get('title', 'AI生成文件')
        ws['A1'] = title
        ws['A1'].font = title_font
        ws.merge_cells('A1:F1')
        
        # 添加標題行
        headers = structure.get('structure', {}).get('headers', ['項目', '描述', '數值'])
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=3, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
        
        # 添加示例數據
        sample_data = [
            ['項目1', '說明1', '100'],
            ['項目2', '說明2', '200'],
            ['項目3', '說明3', '300'],
            ['總計', '', '=SUM(C4:C6)']
        ]
        
        for row_idx, row_data in enumerate(sample_data, 4):
            for col_idx, value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=value)
        
        # 調整列寬
        for column in ws.columns:
            max_length = 0
            column_letter = None
            for cell in column:
                if hasattr(cell, 'column_letter'):
                    column_letter = cell.column_letter
                    if cell.value and len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
            if column_letter and max_length > 0:
                adjusted_width = min(max_length + 2, 30)
                ws.column_dimensions[column_letter].width = adjusted_width
        
        # 保存文件
        import uuid
        filename = f"ai_excel_{uuid.uuid4()}.xlsx"
        file_path = os.path.join('static', 'downloads', filename)
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        wb.save(file_path)
        
        return jsonify({
            'success': True,
            'filename': filename,
            'download_url': f'/download/{filename}',
            'type': 'excel',
            'title': title,
            'description': structure.get('description', ''),
            'file_size': f'{os.path.getsize(file_path) // 1024}KB'
        })
        
    except Exception as e:
        logging.error(f"Excel生成錯誤: {e}")
        return jsonify({'success': False, 'error': f'Excel生成失敗: {str(e)}'})

def create_ai_word_document(structure, images):
    """創建AI生成的Word文件"""
    try:
        if not HAS_DOCX:
            return jsonify({'success': False, 'error': '系統未安裝 python-docx 套件'})
        
        from docx import Document
        from docx.shared import Inches
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        
        doc = Document()
        
        # 添加標題
        title = structure.get('title', 'AI生成文件')
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        # 添加描述
        if structure.get('description'):
            doc.add_paragraph(structure['description'])
            doc.add_paragraph()  # 空行
        
        # 添加內容
        content = structure.get('structure', {}).get('content', '這是AI生成的文件內容。')
        doc.add_paragraph(content)
        
        # 添加圖片
        for i, image in enumerate(images[:2]):  # 最多添加2張圖片
            doc.add_paragraph()  # 空行
            doc.add_heading(f'圖片 {i+1}: {image["description"]}', level=2)
            try:
                doc.add_picture(image['path'], width=Inches(4))
            except Exception as e:
                logging.warning(f"插入圖片失敗: {e}")
                doc.add_paragraph(f"[圖片: {image['description']}]")
        
        # 保存文件
        import uuid
        filename = f"ai_word_{uuid.uuid4()}.docx"
        file_path = os.path.join('static', 'downloads', filename)
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        doc.save(file_path)
        
        return jsonify({
            'success': True,
            'filename': filename,
            'download_url': f'/download/{filename}',
            'type': 'word',
            'title': title,
            'description': structure.get('description', ''),
            'file_size': f'{os.path.getsize(file_path) // 1024}KB',
            'images_generated': images
        })
        
    except Exception as e:
        logging.error(f"Word生成錯誤: {e}")
        return jsonify({'success': False, 'error': f'Word生成失敗: {str(e)}'})

def create_ai_ppt_document(structure, images):
    """創建AI生成的PowerPoint文件"""
    try:
        if not HAS_PPTX:
            return jsonify({'success': False, 'error': '系統未安裝 python-pptx 套件'})
        
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        
        # 標題投影片
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = structure.get('title', 'AI生成簡報')
        subtitle.text = structure.get('description', 'AI智能生成的專業簡報')
        
        # 內容投影片
        slides_data = structure.get('structure', {}).get('slides', [
            {'title': '內容概述', 'content': '這是AI生成的簡報內容', 'image_description': '相關示意圖'}
        ])
        
        for i, slide_data in enumerate(slides_data[:5]):  # 最多5張投影片
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            title = slide.shapes.title
            title.text = slide_data.get('title', f'投影片 {i+1}')
            
            content = slide.placeholders[1]
            content.text = slide_data.get('content', '投影片內容')
            
            # 添加圖片（如果有的話）
            if i < len(images):
                try:
                    slide.shapes.add_picture(images[i]['path'], Inches(7), Inches(2), width=Inches(3))
                except Exception as e:
                    logging.warning(f"PPT插入圖片失敗: {e}")
        
        # 保存文件
        import uuid
        filename = f"ai_ppt_{uuid.uuid4()}.pptx"
        file_path = os.path.join('static', 'downloads', filename)
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        prs.save(file_path)
        
        return jsonify({
            'success': True,
            'filename': filename,
            'download_url': f'/download/{filename}',
            'type': 'powerpoint',
            'title': structure.get('title', 'AI生成簡報'),
            'description': structure.get('description', ''),
            'file_size': f'{os.path.getsize(file_path) // 1024}KB',
            'slides': len(prs.slides),
            'images_generated': images
        })
        
    except Exception as e:
        logging.error(f"PPT生成錯誤: {e}")
        return jsonify({'success': False, 'error': f'PPT生成失敗: {str(e)}'})

def create_google_sheets_template(template_data, language):
    """創建 Google Sheets 範本（實際上創建 CSV 檔案）"""
    try:
        # 創建 CSV 檔案，使用UTF-8-BOM編碼避免中文亂碼
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.csv', mode='w', encoding='utf-8-sig')
        
        # 寫入標題
        title = template_data.get('title', '未命名範本')
        temp_file.write(f"# {title}\n")
        temp_file.write(f"# {template_data.get('description', '')}\n")
        temp_file.write("\n")
        
        # 寫入表頭
        headers = template_data.get('headers', template_data.get('columns', []))
        if headers:
            temp_file.write(','.join(headers) + '\n')
        
        # 寫入資料行
        rows = template_data.get('rows', [])
        for row in rows:
            temp_file.write(','.join([str(cell) for cell in row]) + '\n')
        
        temp_file.close()
        
        # 準備下載檔案
        filename = f"{title}.csv"
        
        # 移動檔案到靜態資料夾（為了下載）
        static_path = os.path.join('static', 'downloads', filename)
        os.makedirs(os.path.dirname(static_path), exist_ok=True)
        shutil.move(temp_file.name, static_path)
        
        return jsonify({
            'success': True,
            'filename': filename,
            'download_url': f'/download/{filename}',
            'message': f'{title} 範本已生成完成，可以匯入 Google Sheets'
        })
        
    except Exception as e:
        logging.error(f"Google Sheets 範本生成錯誤: {e}")
        return jsonify({'success': False, 'error': f'範本生成失敗: {str(e)}'})

def create_spreadsheet(structure, file_type, language):
    """創建試算表"""
    try:
        if not HAS_OPENPYXL:
            return jsonify({'success': False, 'error': '系統未安裝 openpyxl 套件'})
        
        # 創建臨時檔案
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx')
        
        # 使用 openpyxl 創建 Excel 檔案
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = structure.get('title', '工作表1')
        
        # 設置標題樣式
        title_font = Font(size=14, bold=True)
        header_font = Font(size=12, bold=True)
        header_fill = PatternFill(start_color="DDDDDD", end_color="DDDDDD", fill_type="solid")
        
        # 添加標題
        title = structure.get('title')
        if title:
            ws['A1'] = title
            ws['A1'].font = title_font
            ws.merge_cells('A1:F1')
        
        # 添加標題行
        headers = structure.get('headers', structure.get('columns', []))
        if headers:
            start_row = 3 if title else 1
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=start_row, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center')
        
        # 添加數據行
        rows = structure.get('rows', structure.get('data', []))
        if rows:
            start_row = 4 if title else 2
            for row_idx, row in enumerate(rows, start_row):
                for col_idx, value in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
        
        # 自動調整列寬
        for column_cells in ws.columns:
            max_length = 0
            column_letter = None
            for cell in column_cells:
                try:
                    if hasattr(cell, 'column_letter'):
                        column_letter = cell.column_letter
                        if cell.value and len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                except:
                    pass
            if column_letter and max_length > 0:
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
        
        # 保存檔案
        wb.save(temp_file.name)
        temp_file.close()
        
        filename = f"{structure.get('title', 'document')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        
        # 如果是 Google Sheets，嘗試上傳
        if file_type == 'google_sheet':
            try:
                google_url = upload_to_google_sheets(temp_file.name, structure.get('title', 'Document'))
                return jsonify({
                    'success': True,
                    'google_sheet_url': google_url,
                    'title': structure.get('title', 'Document'),
                    'type': 'google_sheet',
                    'description': structure.get('description', '')
                })
            except Exception as e:
                logging.warning(f"Google Sheets 上傳失敗，改為下載檔案: {e}")
        
        # 將檔案保存到靜態目錄以便下載
        import uuid
        file_id = str(uuid.uuid4())
        static_path = f'static/downloads/{file_id}_{filename}'
        os.makedirs('static/downloads', exist_ok=True)
        
        shutil.copy2(temp_file.name, static_path)
        
        # 返回下載資訊
        return jsonify({
            'success': True,
            'download_url': f'/download/{file_id}_{filename}',
            'filename': filename,
            'title': structure.get('title', 'Document'),
            'type': file_type,
            'description': structure.get('description', '')
        })
        
    except Exception as e:
        logging.error(f"試算表創建錯誤: {e}")
        return jsonify({'success': False, 'error': f'試算表創建失敗: {str(e)}'})

def create_word_document(structure, language):
    """創建 Word 文件"""
    try:
        if not HAS_DOCX:
            return jsonify({'success': False, 'error': '系統未安裝 python-docx 套件'})
        
        # 創建臨時檔案
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
        
        # 創建文件
        doc = Document()
        
        # 添加標題
        if 'title' in structure:
            title = doc.add_heading(structure['title'], 0)
            title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER  # 置中
        
        # 添加描述
        if 'description' in structure:
            doc.add_paragraph(structure['description'])
            doc.add_paragraph()  # 空行
        
        # 添加內容
        content = structure.get('content', '')
        if content:
            # 處理內容格式
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    if para.startswith('#'):
                        # 標題
                        level = len(para) - len(para.lstrip('#'))
                        title_text = para.lstrip('# ').strip()
                        doc.add_heading(title_text, level)
                    elif para.startswith('|'):
                        # 表格
                        lines = para.strip().split('\n')
                        if len(lines) > 1:
                            headers = [cell.strip() for cell in lines[0].split('|')[1:-1]]
                            table = doc.add_table(rows=1, cols=len(headers))
                            table.style = 'Table Grid'
                            
                            # 添加標題行
                            hdr_cells = table.rows[0].cells
                            for i, header in enumerate(headers):
                                hdr_cells[i].text = header
                            
                            # 添加數據行
                            for line in lines[2:]:  # 跳過分隔行
                                if line.strip() and not line.strip().startswith('|---'):
                                    cells = [cell.strip() for cell in line.split('|')[1:-1]]
                                    row_cells = table.add_row().cells
                                    for i, cell_value in enumerate(cells):
                                        if i < len(row_cells):
                                            row_cells[i].text = cell_value
                    else:
                        # 普通段落
                        doc.add_paragraph(para)
        
        # 添加sections（如果有）
        sections = structure.get('sections', [])
        for section in sections:
            if isinstance(section, dict):
                if 'title' in section:
                    doc.add_heading(section['title'], 1)
                if 'content' in section:
                    doc.add_paragraph(section['content'])
            else:
                doc.add_paragraph(str(section))
        
        # 保存檔案
        doc.save(temp_file.name)
        temp_file.close()
        
        filename = f"{structure.get('title', 'document')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        
        # 將檔案保存到靜態目錄以便下載
        import uuid
        file_id = str(uuid.uuid4())
        static_path = f'static/downloads/{file_id}_{filename}'
        os.makedirs('static/downloads', exist_ok=True)
        import shutil
        shutil.copy2(temp_file.name, static_path)
        
        # 返回下載資訊
        return jsonify({
            'success': True,
            'download_url': f'/download/{file_id}_{filename}',
            'filename': filename,
            'title': structure.get('title', 'Document'),
            'type': 'word',
            'description': structure.get('description', '')
        })
        
    except Exception as e:
        logging.error(f"Word 文件創建錯誤: {e}")
        return jsonify({'success': False, 'error': f'Word 文件創建失敗: {str(e)}'})

def generate_ppt_document(data):
    """生成PowerPoint文件"""
    try:
        description = data.get('description', '')
        language = data.get('language', 'zh-TW')
        
        if not openai_client.is_configured():
            return jsonify({'success': False, 'error': 'OpenAI API 未配置'})
        
        # 使用AI生成PPT結構
        prompt = f"""
        請根據以下描述生成一個PowerPoint簡報的結構：
        描述：{description}
        語言：{language}
        
        請返回JSON格式，包含以下內容：
        {{
            "title": "簡報標題",
            "slides": [
                {{
                    "title": "投影片標題",
                    "content": "投影片內容（支援換行符號\\n）",
                    "bullets": ["要點1", "要點2", "要點3"]
                }}
            ]
        }}
        請確保投影片數量不超過10張，每張投影片要點不超過5個。
        """
        
        response = openai_client.get_response([{'role': 'user', 'content': prompt}])
        structure = json.loads(response)
        
        # 創建PPT
        prs = Presentation()
        
        # 標題頁
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = structure.get('title', '未命名簡報')
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
        # 內容頁
        for slide_data in structure.get('slides', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # 標題
            title_shape = slide.shapes.title
            title_shape.text = slide_data.get('title', '')
            
            # 內容
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            # 添加內容段落
            content = slide_data.get('content', '')
            if content:
                p = text_frame.paragraphs[0]
                p.text = content
                p.font.size = Pt(18)
            
            # 添加要點
            bullets = slide_data.get('bullets', [])
            for bullet in bullets:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 1
                p.font.size = Pt(16)
        
        # 保存文件
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        temp_file.close()
        
        filename = f"{structure.get('title', 'presentation')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        # 將檔案保存到靜態目錄以便下載
        import uuid
        file_id = str(uuid.uuid4())
        static_path = f'static/downloads/{file_id}_{filename}'
        
        import shutil
        shutil.copy2(temp_file.name, static_path)
        
        # 返回下載資訊
        return jsonify({
            'success': True,
            'download_url': f'/download/{file_id}_{filename}',
            'filename': filename,
            'title': structure.get('title', 'Presentation'),
            'type': 'ppt',
            'description': description,
            'slide_count': len(structure.get('slides', [])) + 1
        })
        
    except Exception as e:
        logging.error(f"PPT生成錯誤: {e}")
        return jsonify({'success': False, 'error': f'PPT生成失敗: {str(e)}'})

def upload_to_google_sheets(file_path, title):
    """上傳到 Google Sheets"""
    # 注意：這需要 Google API 憑證
    # 目前返回模擬連結，實際使用時需要配置 Google API
    logging.info("Google Sheets 功能需要配置 API 憑證")
    raise Exception("Google Sheets 功能需要配置 Google API 憑證")

def get_document_templates(language='zh-TW'):
    """獲取文件範本"""
    if language == 'zh-TW':
        return {
            'spreadsheet': {
                'financial': {
                    'title': '財務報表',
                    'description': '公司財務收支統計表',
                    'headers': ['項目', '預算', '實際', '差異', '百分比', '備註'],
                    'rows': [
                        ['收入', '', '', '', '', ''],
                        ['銷售收入', '100000', '', '', '', ''],
                        ['其他收入', '10000', '', '', '', ''],
                        ['支出', '', '', '', '', ''],
                        ['人事費用', '50000', '', '', '', ''],
                        ['營運費用', '30000', '', '', '', ''],
                        ['行銷費用', '15000', '', '', '', ''],
                        ['淨利潤', '', '', '', '', '']
                    ]
                },
                'project': {
                    'title': '專案計劃',
                    'description': '專案任務規劃與追蹤表',
                    'headers': ['任務', '負責人', '開始日期', '結束日期', '狀態', '進度', '備註'],
                    'rows': [
                        ['需求分析', '', '', '', '未開始', '0%', ''],
                        ['系統設計', '', '', '', '未開始', '0%', ''],
                        ['程式開發', '', '', '', '未開始', '0%', ''],
                        ['測試階段', '', '', '', '未開始', '0%', ''],
                        ['部署上線', '', '', '', '未開始', '0%', '']
                    ]
                },
                'employee': {
                    'title': '員工清單',
                    'description': '公司員工基本資料表',
                    'headers': ['員工編號', '姓名', '部門', '職位', '入職日期', '聯絡電話', '電子郵件'],
                    'rows': [
                        ['', '', '', '', '', '', ''],
                        ['', '', '', '', '', '', ''],
                        ['', '', '', '', '', '', '']
                    ]
                },
                'sales': {
                    'title': '銷售統計',
                    'description': '月度銷售業績統計表',
                    'headers': ['日期', '產品', '數量', '單價', '總額', '業務員', '備註'],
                    'rows': [
                        ['', '', '', '', '', '', ''],
                        ['', '', '', '', '', '', ''],
                        ['', '', '', '', '', '', '']
                    ]
                },
                'inventory': {
                    'title': '庫存管理',
                    'description': '商品庫存追蹤表',
                    'headers': ['商品編號', '商品名稱', '類別', '庫存數量', '安全庫存', '供應商', '最後更新'],
                    'rows': [
                        ['P001', '產品A', '電子產品', '100', '20', '供應商A', '2025-01-01'],
                        ['P002', '產品B', '服飾', '50', '10', '供應商B', '2025-01-01'],
                        ['P003', '產品C', '食品', '200', '30', '供應商C', '2025-01-01']
                    ]
                },
                'hr_attendance': {
                    'title': '員工出勤記錄',
                    'description': '員工每日出勤狀況追蹤表',
                    'headers': ['日期', '員工編號', '姓名', '部門', '上班時間', '下班時間', '工作時數', '遲到', '早退', '備註'],
                    'rows': [
                        ['2025-01-01', 'E001', '張三', '行政部', '09:00', '18:00', '8', '否', '否', ''],
                        ['2025-01-01', 'E002', '李四', '業務部', '09:15', '18:00', '7.75', '是', '否', '遲到15分鐘'],
                        ['2025-01-01', 'E003', '王五', 'IT部', '09:00', '17:30', '7.5', '否', '是', '早退30分鐘'],
                        ['2025-01-02', 'E001', '張三', '行政部', '09:00', '18:00', '8', '否', '否', ''],
                        ['2025-01-02', 'E002', '李四', '業務部', '09:00', '18:00', '8', '否', '否', '']
                    ]
                },
                'hr_salary': {
                    'title': '薪資計算表',
                    'description': '員工月薪計算與明細表',
                    'headers': ['員工編號', '姓名', '部門', '職位', '基本薪資', '加班費', '津貼', '獎金', '應發總額', '勞保', '健保', '所得稅', '應扣總額', '實發金額'],
                    'rows': [
                        ['E001', '張三', '行政部', '經理', '50000', '2000', '3000', '5000', '60000', '1200', '800', '3000', '5000', '55000'],
                        ['E002', '李四', '業務部', '業務員', '40000', '1500', '2000', '8000', '51500', '1000', '600', '2500', '4100', '47400'],
                        ['E003', '王五', 'IT部', '工程師', '60000', '3000', '2500', '10000', '75500', '1500', '1000', '5000', '7500', '68000'],
                        ['E004', '陳六', '財務部', '會計', '45000', '1000', '1500', '3000', '50500', '1100', '700', '2800', '4600', '45900'],
                        ['E005', '劉七', '人事部', '專員', '42000', '800', '1800', '4000', '48600', '1050', '650', '2600', '4300', '44300']
                    ]
                },
                'hr_leave': {
                    'title': '請假申請記錄',
                    'description': '員工請假申請與核准記錄表',
                    'headers': ['申請日期', '員工編號', '姓名', '部門', '請假類型', '開始日期', '結束日期', '請假天數', '請假原因', '核准狀態', '核准人', '備註'],
                    'rows': [
                        ['2025-01-05', 'E001', '張三', '行政部', '事假', '2025-01-10', '2025-01-10', '1', '個人事務', '已核准', '部門主管', ''],
                        ['2025-01-08', 'E002', '李四', '業務部', '病假', '2025-01-12', '2025-01-14', '3', '感冒就醫', '已核准', '部門主管', '需提供診斷證明'],
                        ['2025-01-15', 'E003', '王五', 'IT部', '特休', '2025-01-20', '2025-01-22', '3', '家庭旅遊', '已核准', '部門主管', ''],
                        ['2025-01-18', 'E004', '陳六', '財務部', '喪假', '2025-01-25', '2025-01-27', '3', '親屬過世', '已核准', '人事主管', '法定喪假'],
                        ['2025-01-20', 'E005', '劉七', '人事部', '婚假', '2025-02-01', '2025-02-08', '8', '結婚典禮', '待核准', '部門主管', '需提供結婚證書']
                    ]
                },
                'hr_performance': {
                    'title': '員工績效評估',
                    'description': '員工年度績效考核評估表',
                    'headers': ['員工編號', '姓名', '部門', '職位', '工作品質', '工作效率', '團隊合作', '學習能力', '責任感', '總分', '等級', '評估者', '評估日期'],
                    'rows': [
                        ['E001', '張三', '行政部', '經理', '90', '85', '95', '80', '90', '88', 'A', '總經理', '2024-12-31'],
                        ['E002', '李四', '業務部', '業務員', '85', '90', '80', '85', '85', '85', 'B+', '業務主管', '2024-12-31'],
                        ['E003', '王五', 'IT部', '工程師', '95', '90', '85', '95', '90', '91', 'A', 'IT主管', '2024-12-31'],
                        ['E004', '陳六', '財務部', '會計', '88', '85', '90', '80', '88', '86.2', 'B+', '財務主管', '2024-12-31'],
                        ['E005', '劉七', '人事部', '專員', '80', '85', '95', '90', '85', '87', 'B+', '人事主管', '2024-12-31']
                    ]
                },
                'hr_recruitment': {
                    'title': '招聘管理表',
                    'description': '職位招聘進度追蹤管理表',
                    'headers': ['職位名稱', '部門', '需求人數', '發布日期', '截止日期', '應徵人數', '面試人數', '錄取人數', '招聘狀態', '負責人', '備註'],
                    'rows': [
                        ['軟體工程師', 'IT部', '2', '2025-01-01', '2025-01-31', '25', '8', '1', '進行中', 'IT主管', '需有3年以上經驗'],
                        ['業務專員', '業務部', '3', '2025-01-05', '2025-02-05', '15', '5', '0', '進行中', '業務主管', '具備英文能力佳'],
                        ['會計助理', '財務部', '1', '2025-01-10', '2025-02-10', '12', '3', '0', '進行中', '財務主管', '需會計相關科系'],
                        ['行政助理', '行政部', '1', '2024-12-15', '2025-01-15', '20', '6', '1', '已完成', '行政主管', '已錄取1名'],
                        ['人事專員', '人事部', '1', '2025-01-08', '2025-02-08', '8', '2', '0', '進行中', '人事主管', '需熟悉勞基法']
                    ]
                },
                'hr_training': {
                    'title': '教育訓練記錄',
                    'description': '員工教育訓練參與記錄表',
                    'headers': ['訓練日期', '訓練名稱', '訓練類型', '參與員工', '部門', '訓練時數', '講師', '訓練地點', '考核結果', '證書', '備註'],
                    'rows': [
                        ['2025-01-15', '職場安全衛生', '法規訓練', '全體員工', '全部門', '4', '外部講師', '會議室A', '通過', '是', '年度必修'],
                        ['2025-01-20', 'Excel進階應用', '技能訓練', '行政人員', '行政部', '6', '內部講師', '電腦教室', '通過', '否', ''],
                        ['2025-01-25', '客戶服務技巧', '專業訓練', '業務團隊', '業務部', '8', '外部講師', '訓練中心', '通過', '是', ''],
                        ['2025-02-01', 'Python程式設計', '技術訓練', 'IT人員', 'IT部', '16', '外部講師', '線上課程', '進行中', '否', '為期2天'],
                        ['2025-02-05', '財務分析', '專業訓練', '財務人員', '財務部', '12', '外部講師', '會議室B', '待開始', '是', '']
                    ]
                },
                'budget': {
                    'title': '預算規劃',
                    'description': '年度預算規劃與控制表',
                    'headers': ['預算項目', '預算金額', '已使用', '剩餘預算', '使用率', '負責部門', '備註'],
                    'rows': [
                        ['人事費用', '2400000', '200000', '2200000', '8.3%', '人事部', '含薪資、獎金、保險'],
                        ['辦公用品', '120000', '25000', '95000', '20.8%', '行政部', ''],
                        ['行銷推廣', '500000', '150000', '350000', '30%', '行銷部', '包含廣告、活動費用'],
                        ['設備採購', '800000', '320000', '480000', '40%', 'IT部', '電腦、軟體、網路設備'],
                        ['差旅費用', '200000', '35000', '165000', '17.5%', '全部門', '']
                    ]
                },
                'customer': {
                    'title': '客戶資料',
                    'description': '客戶基本資料管理表',
                    'headers': ['客戶編號', '公司名稱', '聯絡人', '職稱', '電話', '電子郵件', '地址', '客戶類型', '業務負責', '備註'],
                    'rows': [
                        ['C001', 'ABC科技公司', '王經理', '採購經理', '02-12345678', 'wang@abc.com', '台北市信義區', 'A級客戶', '李四', '長期合作夥伴'],
                        ['C002', 'XYZ貿易有限公司', '陳總監', '營運總監', '03-87654321', 'chen@xyz.com', '桃園市中壢區', 'B級客戶', '李四', ''],
                        ['C003', '123製造股份有限公司', '林主任', '品管主任', '04-11223344', 'lin@123.com', '台中市西屯區', 'A級客戶', '張五', '品質要求嚴格'],
                        ['C004', '456服務集團', '黃副理', '業務副理', '07-99887766', 'huang@456.com', '高雄市前鎮區', 'C級客戶', '張五', '新客戶'],
                        ['C005', '789科技股份有限公司', '吳總經理', '總經理', '02-55443322', 'wu@789.com', '新北市板橋區', 'A級客戶', '李四', 'VIP客戶']
                    ]
                },
                'grades': {
                    'title': '學生成績',
                    'description': '學生學期成績記錄表',
                    'headers': ['學號', '姓名', '班級', '國文', '英文', '數學', '自然', '社會', '總分', '平均', '排名'],
                    'rows': [
                        ['A001', '小明', '三年甲班', '85', '90', '88', '92', '87', '442', '88.4', '3'],
                        ['A002', '小華', '三年甲班', '92', '85', '95', '89', '91', '452', '90.4', '1'],
                        ['A003', '小美', '三年甲班', '88', '92', '85', '90', '89', '444', '88.8', '2'],
                        ['A004', '小強', '三年甲班', '80', '78', '82', '85', '83', '408', '81.6', '5'],
                        ['A005', '小麗', '三年甲班', '90', '88', '87', '86', '85', '436', '87.2', '4']
                    ]
                },
                'schedule': {
                    'title': '工作排程',
                    'description': '每週工作時程安排表',
                    'headers': ['時間', '星期一', '星期二', '星期三', '星期四', '星期五', '星期六', '星期日'],
                    'rows': [
                        ['09:00-10:00', '', '', '', '', '', '', ''],
                        ['10:00-11:00', '', '', '', '', '', '', ''],
                        ['11:00-12:00', '', '', '', '', '', '', ''],
                        ['13:00-14:00', '', '', '', '', '', '', ''],
                        ['14:00-15:00', '', '', '', '', '', '', ''],
                        ['15:00-16:00', '', '', '', '', '', '', ''],
                        ['16:00-17:00', '', '', '', '', '', '', '']
                    ]
                },
                'budget': {
                    'title': '年度預算',
                    'description': '年度預算規劃與執行表',
                    'headers': ['項目', 'Q1預算', 'Q1實際', 'Q2預算', 'Q2實際', 'Q3預算', 'Q3實際', 'Q4預算', 'Q4實際', '年度總計'],
                    'rows': [
                        ['營收', '250000', '', '280000', '', '300000', '', '320000', '', ''],
                        ['人事成本', '100000', '', '105000', '', '110000', '', '115000', '', ''],
                        ['營運成本', '80000', '', '85000', '', '90000', '', '95000', '', ''],
                        ['行銷費用', '30000', '', '35000', '', '40000', '', '45000', '', ''],
                        ['淨利潤', '40000', '', '55000', '', '60000', '', '65000', '', '']
                    ]
                },
                'attendance': {
                    'title': '出勤記錄',
                    'description': '員工出勤統計表',
                    'headers': ['員工姓名', '部門', '工號', '出勤天數', '遲到次數', '早退次數', '請假天數', '加班時數', '備註'],
                    'rows': [
                        ['', '', '', '', '', '', '', '', ''],
                        ['', '', '', '', '', '', '', '', ''],
                        ['', '', '', '', '', '', '', '', '']
                    ]
                }
            },
            'document': {
                'meeting': {
                    'title': '會議紀錄',
                    'description': '標準會議紀錄範本',
                    'content': '''會議紀錄

會議主題：
日期時間：
地點：
主持人：
出席人員：

## 議程
1. 
2. 
3. 

## 討論內容


## 決議事項
1. 
2. 
3. 

## 待辦事項
| 項目 | 負責人 | 期限 |
|------|--------|------|
|      |        |      |

## 下次會議
日期：
地點：'''
                },
                'proposal': {
                    'title': '項目提案',
                    'description': '標準項目提案書範本',
                    'content': '''項目提案書

提案標題：
提案人：
提案日期：

## 一、背景與目標
### 1.1 專案背景

### 1.2 專案目標

## 二、解決方案
### 2.1 方案概述

### 2.2 實施計劃

## 三、資源需求
### 3.1 人力資源

### 3.2 預算需求

## 四、效益評估
### 4.1 預期效益

### 4.2 風險評估

## 五、時程規劃
| 階段 | 時間 | 里程碑 |
|------|------|--------|
|      |      |        |

## 六、結論'''
                },
                'contract': {
                    'title': '合約範本',
                    'description': '標準合約協議書範本',
                    'content': '''合約協議書

甲方：
乙方：
簽約日期：

## 第一條 合約標的
本合約內容為：

## 第二條 權利義務
### 甲方權利義務：
1. 
2. 
3. 

### 乙方權利義務：
1. 
2. 
3. 

## 第三條 付款條件
付款方式：
付款期限：
付款金額：

## 第四條 違約責任
1. 
2. 
3. 

## 第五條 爭議解決
如有爭議，雙方應協商解決，協商不成時，提交仲裁委員會仲裁。

## 第六條 其他條款
1. 本合約一式兩份，甲乙雙方各執一份
2. 本合約自雙方簽字蓋章之日起生效
3. 未盡事宜，雙方可另行協商

甲方簽名：_____________ 日期：_________
乙方簽名：_____________ 日期：_________'''
                },
                'report': {
                    'title': '工作報告',
                    'description': '工作成果報告範本',
                    'content': '''工作報告

報告人：
報告日期：
報告期間：

## 一、工作概況
### 1.1 主要工作內容
-
-
-

### 1.2 完成情況
-
-
-

## 二、具體成果
### 2.1 量化指標
| 指標 | 目標值 | 實際值 | 達成率 |
|------|--------|--------|--------|
|      |        |        |        |

### 2.2 質化成果
-
-
-

## 三、遇到的問題
### 3.1 主要困難
-
-
-

### 3.2 解決方案
-
-
-

## 四、下期工作計劃
### 4.1 工作重點
-
-
-

### 4.2 預期目標
-
-
-

## 五、建議與反思
-
-
-'''
                },
                'policy': {
                    'title': '公司政策',
                    'description': '公司制度政策範本',
                    'content': '''公司政策文件

政策名稱：
制定部門：
生效日期：
版本：1.0

## 一、政策目的
本政策旨在：

## 二、適用範圍
本政策適用於：

## 三、定義
-
-
-

## 四、政策內容
### 4.1 基本原則
-
-
-

### 4.2 具體規定
-
-
-

### 4.3 執行程序
1. 
2. 
3. 

## 五、責任分工
| 職位 | 責任 |
|------|------|
|      |      |

## 六、監督檢查
-
-
-

## 七、違規處理
-
-
-

## 八、附則
1. 本政策解釋權歸人事部門
2. 本政策自發布之日起執行
3. 如有修訂，另行通知

制定人：_____________ 日期：_________
審核人：_____________ 日期：_________
批准人：_____________ 日期：_________'''
                },
                'resume': {
                    'title': '履歷表',
                    'description': '個人履歷表範本',
                    'content': '''個人履歷

# 基本資料
姓名：
性別：
出生年月日：
聯絡電話：
電子郵件：
現居地址：

# 求職意向
應徵職位：
期望薪資：
可到職日：

# 學歷
## 最高學歷
學校名稱：
科系：
學位：
畢業年月：
成績：

## 其他學歷

# 工作經驗
## 目前/最近工作
公司名稱：
職位名稱：
任職期間：
工作內容：
-
-
-

## 前一份工作
公司名稱：
職位名稱：
任職期間：
工作內容：
-
-
-

# 專業技能
## 語言能力
- 中文：母語
- 英文：
- 其他：

## 電腦技能
- 辦公室軟體：
- 程式語言：
- 其他技能：

# 證照/證書
- 
- 
- 

# 自傳
請簡述個人特質、求學/工作經驗、未來規劃等。

# 其他資訊
興趣嗜好：
參與社團：
獲獎記錄：'''
                },
                'technical': {
                    'title': '技術文檔',
                    'description': '技術說明文件範本',
                    'content': '''技術文檔

文檔標題：
版本：1.0
作者：
最後更新：

# 概述
## 目的
本文檔的目的是...

## 適用範圍
-
-
-

# 系統架構
## 整體架構

## 主要組件
### 組件一
功能說明：
技術規格：

### 組件二
功能說明：
技術規格：

# 安裝說明
## 系統需求
- 作業系統：
- 硬體需求：
- 軟體需求：

## 安裝步驟
1. 
2. 
3. 

# 使用說明
## 基本操作
### 功能一
操作步驟：
1. 
2. 
3. 

### 功能二
操作步驟：
1. 
2. 
3. 

# API 說明
## 端點列表
| 端點 | 方法 | 描述 | 參數 |
|------|------|------|------|
|      |      |      |      |

## 範例
### 請求範例
```
API 請求範例
```

### 回應範例
```
API 回應範例
```

# 故障排除
## 常見問題
### 問題一
現象：
解決方案：

### 問題二
現象：
解決方案：

# 版本歷史
| 版本 | 日期 | 修改內容 | 作者 |
|------|------|----------|------|
| 1.0  |      |          |      |'''
                },
                'report': {
                    'title': '報告書',
                    'description': '標準報告書範本',
                    'content': '''報告書

報告標題：
報告人：
報告日期：
報告對象：

# 摘要
本報告簡述...

# 一、前言
## 1.1 報告背景

## 1.2 報告目的

## 1.3 報告範圍

# 二、方法論
## 2.1 研究方法

## 2.2 資料來源

## 2.3 分析工具

# 三、發現與分析
## 3.1 主要發現
-
-
-

## 3.2 數據分析
### 分析一
結果：
說明：

### 分析二
結果：
說明：

# 四、結論
## 4.1 總結

## 4.2 限制

# 五、建議
## 5.1 短期建議
1. 
2. 
3. 

## 5.2 長期建議
1. 
2. 
3. 

# 六、附錄
## 附錄A：數據表格

## 附錄B：參考資料

---
報告完成日期：
報告人簽名：____________'''
                },
                'business_plan': {
                    'title': '商業計劃',
                    'description': '商業計劃書範本',
                    'content': '''商業計劃書

公司名稱：
計劃日期：
計劃期間：

# 執行摘要
## 公司概述

## 產品/服務概述

## 市場機會

## 競爭優勢

## 財務預測

## 資金需求

# 一、公司描述
## 1.1 公司使命

## 1.2 公司願景

## 1.3 核心價值

## 1.4 公司歷史

## 1.5 法律結構

# 二、市場分析
## 2.1 行業概況

## 2.2 目標市場

## 2.3 市場規模

## 2.4 市場趨勢

# 三、競爭分析
## 3.1 競爭對手分析
| 競爭對手 | 優勢 | 劣勢 | 市場佔有率 |
|----------|------|------|------------|
|          |      |      |            |

## 3.2 競爭策略

# 四、產品與服務
## 4.1 產品/服務描述

## 4.2 特色與優勢

## 4.3 開發階段

## 4.4 智慧財產權

# 五、行銷與銷售策略
## 5.1 行銷策略

## 5.2 銷售策略

## 5.3 定價策略

## 5.4 通路策略

# 六、營運計劃
## 6.1 營運流程

## 6.2 設施與設備

## 6.3 技術需求

# 七、管理團隊
## 7.1 組織架構

## 7.2 核心團隊
### 執行長
姓名：
經歷：
職責：

### 其他主管

## 7.3 人力資源計劃

# 八、財務計劃
## 8.1 資金需求

## 8.2 財務預測
### 損益預測
| 年度 | 營收 | 支出 | 淨利 |
|------|------|------|------|
| 第1年|      |      |      |
| 第2年|      |      |      |
| 第3年|      |      |      |

### 現金流預測

## 8.3 投資回報

# 九、風險分析
## 9.1 市場風險

## 9.2 技術風險

## 9.3 財務風險

## 9.4 風險控制措施

# 十、附錄
## 附錄A：財務詳細預測

## 附錄B：市場研究數據

## 附錄C：技術規格'''
                },
                'product_desc': {
                    'title': '產品說明',
                    'description': '產品說明書範本',
                    'content': '''產品說明書

產品名稱：
產品型號：
版本：
發布日期：

# 產品概述
## 產品簡介

## 主要特色
-
-
-

## 適用對象

# 產品規格
## 技術規格
| 項目 | 規格 |
|------|------|
| 尺寸 |      |
| 重量 |      |
| 材質 |      |
| 顏色 |      |

## 系統需求
- 作業系統：
- 硬體需求：
- 軟體需求：

# 功能說明
## 主要功能
### 功能一
描述：
使用方法：

### 功能二
描述：
使用方法：

## 進階功能
### 功能三
描述：
使用方法：

# 使用指南
## 快速入門
1. 
2. 
3. 

## 詳細操作
### 設定步驟
1. 
2. 
3. 

### 常用操作
1. 
2. 
3. 

# 安全須知
## 使用注意事項
-
-
-

## 維護保養
-
-
-

# 故障排除
## 常見問題
### Q1: 
A1: 

### Q2: 
A2: 

## 聯絡資訊
客服電話：
電子郵件：
官方網站：

# 保固資訊
保固期間：
保固範圍：
保固條件：

# 技術支援
支援時間：
聯絡方式：
線上資源：'''
                },
                'manual': {
                    'title': '操作手冊',
                    'description': '標準操作手冊範本',
                    'content': '''操作手冊

手冊名稱：
版本：1.0
制定日期：
最後更新：

## 一、概述
### 1.1 目的
本手冊旨在：

### 1.2 適用範圍
-
-
-

## 二、準備工作
### 2.1 所需工具
-
-
-

### 2.2 前置條件
-
-
-

## 三、操作步驟
### 3.1 步驟一：
1. 
2. 
3. 

### 3.2 步驟二：
1. 
2. 
3. 

### 3.3 步驟三：
1. 
2. 
3. 

## 四、注意事項
### 4.1 安全須知
-
-
-

### 4.2 重要提醒
-
-
-

## 五、常見問題
### Q1: 
A1: 

### Q2: 
A2: 

### Q3: 
A3: 

## 六、聯絡資訊
技術支援：
緊急聯絡：
更新維護：'''
                }
            }
        }
    
    # 英文版本
    return {
        'spreadsheet': {
            'financial': {
                'title': 'Financial Report',
                'description': 'Company financial income and expense statement',
                'headers': ['Item', 'Budget', 'Actual', 'Variance', 'Percentage', 'Notes'],
                'rows': [
                    ['Income', '', '', '', '', ''],
                    ['Sales Revenue', '100000', '', '', '', ''],
                    ['Other Income', '10000', '', '', '', ''],
                    ['Expenses', '', '', '', '', ''],
                    ['Personnel Costs', '50000', '', '', '', ''],
                    ['Operating Expenses', '30000', '', '', '', ''],
                    ['Marketing Costs', '15000', '', '', '', ''],
                    ['Net Profit', '', '', '', '', '']
                ]
            }
        },
        'document': {
            'meeting': {
                'title': 'Meeting Minutes',
                'description': 'Standard meeting minutes template',
                'content': '''Meeting Minutes

Meeting Topic:
Date & Time:
Location:
Chairperson:
Attendees:

## Agenda
1. 
2. 
3. 

## Discussion


## Decisions
1. 
2. 
3. 

## Action Items
| Item | Responsible | Deadline |
|------|-------------|----------|
|      |             |          |

## Next Meeting
Date:
Location:'''
            }
        }
    }

@app.route('/clear_history', methods=['POST'])
def clear_history():
    """Clear chat history"""
    try:
        session.pop('chat_history', None)
        return jsonify({'success': True, 'message': 'Chat history cleared'})
    except Exception as e:
        logging.error(f"Error clearing history: {str(e)}")
        return jsonify({'error': 'Failed to clear history', 'success': False}), 500

@app.route('/generate_qr', methods=['POST'])
def generate_qr():
    """Generate QR code"""
    try:
        data = request.json
        text = data.get('text', '')
        size = data.get('size', 10)
        border = data.get('border', 4)
        
        if not text:
            return jsonify({'error': 'Text is required'}), 400
        
        # 創建QR碼
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=size,
            border=border,
        )
        qr.add_data(text)
        qr.make(fit=True)
        
        # 生成圖片
        img = qr.make_image(fill_color="black", back_color="white")
        
        # 轉換為base64
        buffer = BytesIO()
        img.save(buffer, format='PNG')
        buffer.seek(0)
        
        img_base64 = base64.b64encode(buffer.getvalue()).decode()
        
        return jsonify({
            'success': True,
            'image': f'data:image/png;base64,{img_base64}',
            'download_data': img_base64
        })
        
    except Exception as e:
        logging.error(f"Error generating QR code: {str(e)}")
        return jsonify({'error': 'Failed to generate QR code'}), 500

@app.route('/generate_barcode', methods=['POST'])
def generate_barcode():
    """Generate barcode"""
    try:
        data = request.json
        text = data.get('text', '')
        barcode_type = data.get('type', 'code128')
        
        if not text:
            return jsonify({'error': 'Text is required'}), 400
        
        # 支援的條碼類型
        barcode_types = {
            'code128': Code128,
            'code39': Code39,
            'ean13': EAN13,
            'ean8': EAN8,
            'upc': UPCA,
            'isbn13': ISBN13,
            'isbn10': ISBN10,
            'issn': ISSN
        }
        
        if barcode_type not in barcode_types:
            return jsonify({'error': 'Invalid barcode type'}), 400
        
        # 生成條碼
        barcode_class = barcode_types[barcode_type]
        code = barcode_class(text, writer=ImageWriter())
        
        buffer = BytesIO()
        code.write(buffer)
        buffer.seek(0)
        
        img_base64 = base64.b64encode(buffer.getvalue()).decode()
        
        return jsonify({
            'success': True,
            'image': f'data:image/png;base64,{img_base64}',
            'download_data': img_base64,
            'type': barcode_type
        })
        
    except Exception as e:
        logging.error(f"Error generating barcode: {str(e)}")
        return jsonify({'error': f'Failed to generate barcode: {str(e)}'}), 500



@app.route('/download/<filename>')
def download_file(filename):
    """通用文件下載端點"""
    try:
        file_path = f'static/downloads/{filename}'
        if os.path.exists(file_path):
            return send_file(file_path, as_attachment=True)
        else:
            return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        logging.error(f"Download error: {str(e)}")
        return jsonify({'error': 'Download failed'}), 500

@app.route('/mahjong/start', methods=['POST'])
def start_mahjong_game():
    """開始麻將遊戲"""
    try:
        # 創建完整的144張牌組
        base_tiles = [
            '1萬', '2萬', '3萬', '4萬', '5萬', '6萬', '7萬', '8萬', '9萬',
            '1筒', '2筒', '3筒', '4筒', '5筒', '6筒', '7筒', '8筒', '9筒',
            '1條', '2條', '3條', '4條', '5條', '6條', '7條', '8條', '9條',
            '東', '南', '西', '北', '中', '發', '白'
        ]
        
        # 每種牌4張
        deck = []
        for tile in base_tiles:
            for _ in range(4):
                deck.append(tile)
        
        # 洗牌
        import random
        random.shuffle(deck)
        
        # 發牌給玩家和3個電腦
        player_hand = sorted(deck[:13])
        computer_hands = [deck[13:26], deck[26:39], deck[39:52]]
        remaining_deck = deck[52:]
        
        # 儲存遊戲狀態到session
        session['mahjong_game'] = {
            'player_hand': player_hand,
            'computer_hands': computer_hands,
            'deck': remaining_deck,
            'discarded_tiles': [],
            'player_melded': [],
            'computer_melded': [[], [], []],
            'current_player': 0,  # 0=玩家, 1-3=電腦
            'last_discarded_tile': None,
            'last_discarded_player': None,
            'game_started': True
        }
        
        return jsonify({
            'success': True,
            'player_hand': player_hand,
            'hand_count': [len(computer_hands[0]), len(computer_hands[1]), len(computer_hands[2])],
            'remaining_tiles': len(remaining_deck)
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/mahjong/discard', methods=['POST'])
def discard_mahjong_tile():
    """打出麻將牌"""
    try:
        data = request.get_json()
        tile_index = data.get('tile_index')
        
        game_state = session.get('mahjong_game')
        if not game_state:
            return jsonify({'success': False, 'error': '遊戲未開始'})
        
        if tile_index < 0 or tile_index >= len(game_state['player_hand']):
            return jsonify({'success': False, 'error': '無效的牌索引'})
        
        # 玩家打牌
        discarded_tile = game_state['player_hand'].pop(tile_index)
        game_state['discarded_tiles'].append(discarded_tile)
        game_state['last_discarded_tile'] = discarded_tile
        game_state['last_discarded_player'] = 0
        game_state['current_player'] = 1
        
        session['mahjong_game'] = game_state
        
        return jsonify({
            'success': True,
            'discarded_tile': discarded_tile,
            'player_hand': game_state['player_hand'],
            'discarded_tiles': game_state['discarded_tiles']
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/mahjong/draw', methods=['POST'])
def draw_mahjong_tile():
    """摸牌"""
    try:
        game_state = session.get('mahjong_game')
        if not game_state or not game_state['deck']:
            return jsonify({'success': False, 'error': '無法摸牌'})
        
        new_tile = game_state['deck'].pop()
        game_state['player_hand'].append(new_tile)
        game_state['player_hand'].sort()
        
        session['mahjong_game'] = game_state
        
        return jsonify({
            'success': True,
            'new_tile': new_tile,
            'player_hand': game_state['player_hand'],
            'remaining_tiles': len(game_state['deck'])
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/mahjong/action', methods=['POST'])
def execute_mahjong_action():
    """執行麻將動作（吃、碰、槓、胡）"""
    try:
        data = request.get_json()
        action = data.get('action')
        
        game_state = session.get('mahjong_game')
        if not game_state:
            return jsonify({'success': False, 'error': '遊戲未開始'})
        
        if action == 'hu':
            return jsonify({
                'success': True,
                'action': 'hu',
                'message': '恭喜胡牌！',
                'game_over': True
            })
        
        elif action == 'pong':
            tile = game_state['last_discarded_tile']
            if tile and game_state['player_hand'].count(tile) >= 2:
                # 執行碰牌
                for _ in range(2):
                    game_state['player_hand'].remove(tile)
                game_state['player_melded'].append([tile, tile, tile])
                if game_state['discarded_tiles']:
                    game_state['discarded_tiles'].pop()
                
                session['mahjong_game'] = game_state
                
                return jsonify({
                    'success': True,
                    'action': 'pong',
                    'player_hand': game_state['player_hand'],
                    'player_melded': game_state['player_melded'],
                    'discarded_tiles': game_state['discarded_tiles']
                })
        
        return jsonify({'success': False, 'error': '無法執行該動作'})
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/mahjong/computer_turn', methods=['POST'])
def computer_mahjong_turn():
    """電腦回合"""
    try:
        game_state = session.get('mahjong_game')
        if not game_state:
            return jsonify({'success': False, 'error': '遊戲未開始'})
        
        current_player = game_state['current_player']
        if current_player < 1 or current_player > 3:
            current_player = 1
            game_state['current_player'] = 1
        
        computer_index = current_player - 1
        computer_hand = game_state['computer_hands'][computer_index]
        
        # 電腦摸牌
        if game_state['deck']:
            new_tile = game_state['deck'].pop()
            computer_hand.append(new_tile)
        
        # 電腦打牌（簡單AI）
        discarded_tile = None
        if computer_hand:
            import random
            discard_index = random.randint(0, len(computer_hand) - 1)
            discarded_tile = computer_hand.pop(discard_index)
            game_state['discarded_tiles'].append(discarded_tile)
            game_state['last_discarded_tile'] = discarded_tile
            game_state['last_discarded_player'] = current_player
        
        # 下一個玩家
        game_state['current_player'] = (current_player + 1) % 4
        
        session['mahjong_game'] = game_state
        
        # 檢查玩家是否可以執行動作
        available_actions = []
        if discarded_tile:
            player_hand = game_state['player_hand']
            if can_chi(player_hand, discarded_tile):
                available_actions.append('吃')
            if can_pong(player_hand, discarded_tile):
                available_actions.append('碰')
            if can_hu(player_hand, discarded_tile):
                available_actions.append('胡')
        
        return jsonify({
            'success': True,
            'discarded_tile': discarded_tile or '',
            'computer_player': current_player,
            'discarded_tiles': game_state['discarded_tiles'],
            'hand_counts': [len(h) for h in game_state['computer_hands']],
            'current_player': game_state['current_player'],
            'available_actions': available_actions
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/mahjong')
def mahjong_page():
    """麻將遊戲頁面"""
    return render_template('mahjong.html')

@app.route('/health')
def health():
    """Health check endpoint"""
    return jsonify({'status': 'healthy', 'openai_configured': openai_client.is_configured()})

# 麻將遊戲邏輯函數
def can_chi(hand, tile):
    """檢查是否可以吃牌"""
    if tile in ['東', '南', '西', '北', '中', '發', '白']:
        return False
    
    try:
        num = int(tile[0])
        suit = tile[1]
    except:
        return False
    
    hand_tiles = [t for t in hand if t.endswith(suit)]
    hand_nums = []
    for t in hand_tiles:
        try:
            hand_nums.append(int(t[0]))
        except:
            continue
    
    # 檢查順子可能性
    if (num-2 in hand_nums and num-1 in hand_nums) or \
       (num-1 in hand_nums and num+1 in hand_nums) or \
       (num+1 in hand_nums and num+2 in hand_nums):
        return True
    
    return False

def can_pong(hand, tile):
    """檢查是否可以碰牌"""
    return hand.count(tile) >= 2

def can_hu(hand, tile=None):
    """檢查是否可以胡牌（簡化版）"""
    test_hand = hand.copy()
    if tile:
        test_hand.append(tile)
    
    if len(test_hand) != 14:
        return False
    
    # 簡化的胡牌判斷：檢查對子和組合
    for unique_tile in set(test_hand):
        if test_hand.count(unique_tile) >= 2:
            temp_hand = test_hand.copy()
            temp_hand.remove(unique_tile)
            temp_hand.remove(unique_tile)
            
            if check_complete_hand(temp_hand):
                return True
    
    return False

def check_complete_hand(hand):
    """檢查手牌是否完整"""
    if len(hand) == 0:
        return True
    
    if len(hand) % 3 != 0:
        return False
    
    # 檢查刻子
    for unique_tile in set(hand):
        if hand.count(unique_tile) >= 3:
            temp_hand = hand.copy()
            for _ in range(3):
                temp_hand.remove(unique_tile)
            return check_complete_hand(temp_hand)
    
    return False

if __name__ == '__main__':
    app.run(host="0.0.0.0", port=5000, debug=True)
